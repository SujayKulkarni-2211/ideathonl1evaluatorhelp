import google.generativeai as genai
import pptx
from config import GEMINI_API_KEY

genai.configure(api_key=GEMINI_API_KEY)

def extract_text_and_images(ppt_path):
    prs = pptx.Presentation(ppt_path)
    extracted_text = []
    image_count = 0

    for slide in prs.slides:
        slide_text = " ".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
        extracted_text.append(slide_text)

        for shape in slide.shapes:
            if hasattr(shape, "image"):
                image_count += 1

    return " ".join(extracted_text), image_count

def evaluate_ppt(ppt_path):
    extracted_text, image_count = extract_text_and_images(ppt_path)

    prompt = f"""
    You are evaluating an Ideathon presentation based on multiple weighted criteria.  
    **Each category has a specific maximum score. You must provide integer scores out of the given weightage.**  

    **Scoring Criteria & Weightage:**  
    1. **Relevance to Problem Statement (out of 20)**
       - How well does the idea align with the problem statement?  
       - Example: A chatbot for customer service scores **low** in a cybersecurity challenge.  
    2. **Innovation & Creativity (out of 20)**  
       - Is the approach unique? Does it introduce fresh thinking?  
       - Example: A blockchain-based cybersecurity protocol would score **high**.  
    3. **Usefulness & Practicality (out of 15)**  
       - Can the solution be applied in the real world? Does it solve a problem effectively?  
       - Example: An AI-driven fraud detection model for banking would score **high**.  
    4. **Originality of Idea (out of 15)**  
       - Does the idea stand out? Is it different from existing solutions?  
       - Example: A **new ML model** trained on domain-specific datasets scores **high**, but using a standard model scores **low**.  
    5. **Technical Feasibility (out of 10)**  
       - Can the idea be implemented with current technology?  
       - Example: A **quantum cryptography solution for IoT** is futuristic, so it would score **low**.  
    6. **Future Scope & Market Potential (out of 10)**  
       - Does the solution have long-term potential? Can it scale in the market?  
       - Example: A **subscription-based AI security tool** has more market potential than a one-time downloadable script.  
    7. **Sustainability & Green Impact (out of 5)**  
       - Does the idea consider environmental or ethical impacts?  
       - Example: A **low-energy AI model** scores **higher** than a solution with excessive power consumption.  

    **Instructions for Scoring:**  
    - Provide scores in this exact format (only numbers, comma-separated):  
      **Relevance, Innovation, Usefulness, Originality, Feasibility, Future Scope, Sustainability**  
    - Example Output: **18, 16, 12, 14, 8, 7, 4**  

    **Now, evaluate the following presentation:**  
    {extracted_text}
    """

    model = genai.GenerativeModel("gemini-2.0-flash")
    response = model.generate_content(prompt)

    # Ensure AI outputs the expected format
    scores_list = [int(score.strip()) for score in response.text.strip().split(",")]

    # Flagging logic based on individual category cutoffs
    flag = "Red" if any(score < (weight * 0.6) for score, weight in zip(scores_list, [20, 20, 15, 15, 10, 10, 5])) else \
           "Orange" if image_count > 5 else "Green"

    return scores_list, flag
